import os
import anthropic
from pptx import Presentation
from database import get_decks_by_category, get_conn

QUALITY_THRESHOLD = 4  # decks scoring below this are removed

CATEGORY_CRITERIA = {
    "rtb": "RTB (Reason to Believe) deck — should have clear audience insight, data-backed targeting rationale, and professional structure",
    "eoc": "EOC (End of Campaign) report — should have campaign results, performance data, clear narrative, and actionable insights",
    "qbr": "QBR (Quarterly Business Review) — should have strategic overview, performance trends, and forward-looking recommendations",
    "intro": "Broadlab introduction deck — should clearly introduce Broadlab's proposition, capabilities, and be polished for new client meetings",
    "optimisation": "Campaign optimisation report — should have mid-campaign performance data, what's working, and recommended changes",
}


def extract_text(pptx_path):
    try:
        prs = Presentation(pptx_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                slides.append(f"Slide {i+1}: " + " | ".join(texts[:4]))
        return "\n".join(slides[:12])  # first 12 slides max
    except Exception:
        return ""


def score_deck(client, category, filename, content):
    criteria = CATEGORY_CRITERIA.get(category, "professional client-facing presentation")
    prompt = f"""You are reviewing a Broadlab {category.upper()} presentation to decide if it belongs in a quality ranking.

Deck name: {filename}
Type: {criteria}

Slide content:
{content}

Score this deck 1-7 based on:
- 1-2: Template, placeholder, unfinished draft, or clearly low quality
- 3-4: Real deck but weak structure, missing data, or generic content
- 5-6: Solid client-facing deck with good structure and content
- 7: Exceptional — clear, compelling, well-structured

Respond with ONLY a JSON object like this:
{{"score": 5, "reason": "one sentence"}}"""

    message = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=100,
        messages=[{"role": "user", "content": prompt}]
    )

    import json
    text = message.content[0].text.strip()
    # Extract JSON even if there's surrounding text
    start = text.find('{')
    end = text.rfind('}') + 1
    if start >= 0 and end > start:
        result = json.loads(text[start:end])
        return result.get('score', 5), result.get('reason', '')
    return 5, "parse error"


def analyse_category(category, api_key):
    client = anthropic.Anthropic(api_key=api_key)
    decks = get_decks_by_category(category)

    results = []
    removed = 0

    for deck in decks:
        content = extract_text(deck['path'])
        if not content:
            results.append({'id': deck['id'], 'filename': deck['filename'],
                           'score': None, 'reason': 'Could not read file', 'kept': False})
            removed += 1
            continue

        score, reason = score_deck(client, category, deck['filename'], content)
        keep = score >= QUALITY_THRESHOLD

        if not keep:
            with get_conn() as conn:
                conn.execute("DELETE FROM decks WHERE id = ?", (deck['id'],))
            removed += 1

        results.append({
            'id': deck['id'],
            'filename': deck['filename'],
            'score': score,
            'reason': reason,
            'kept': keep
        })

    return results, removed
